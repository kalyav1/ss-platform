# =========================================================
# universal_normalizer.py
# =========================================================

import os
import re
import json
import logging
from pathlib import Path

import yaml
import pandas as pd

from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

from PIL import Image
import pytesseract


# =========================================================
# LOGGING CONFIG
# =========================================================

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s"
)

logger = logging.getLogger("universal_normalizer")


# =========================================================
# CONFIGURATION
# =========================================================

CONFIG = {
    "MAX_FILE_SIZE_MB": 20,
    "MAX_INPUTS": 25,
    "MAX_CONTENT_CHARS_PER_ITEM": 50000,
    "MAX_FINAL_CONTEXT_CHARS": 200000,
    "SUPPORTED_EXTENSIONS": {
        ".txt", ".md", ".json", ".yaml", ".yml",
        ".csv", ".xlsx", ".docx", ".pptx", ".pdf",
        ".png", ".jpg", ".jpeg"
    }
}


# =========================================================
# TEXT NORMALIZATION
# =========================================================

def normalize_text(text):

    if not text:
        return ""

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.encode("utf-8", errors="ignore").decode("utf-8")
    text = text.replace("\t", " ")
    text = re.sub(r"[ ]{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def truncate_text(text, max_chars):
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


# =========================================================
# SIZE CHECK
# =========================================================

def size_limit_check(path):
    file_size_mb = os.path.getsize(path) / (1024 * 1024)
    return file_size_mb <= CONFIG["MAX_FILE_SIZE_MB"]


# =========================================================
# CLASSIFICATION
# =========================================================

def classify_input(item):
    if isinstance(item, str) and os.path.exists(item):
        return "FILE"
    return "RAW_TEXT"


# =========================================================
# FILE EXTRACTORS
# =========================================================

def extract_text_file(path):
    logger.info(f"Extracting TXT/MD: {path}")
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_json(path):
    logger.info(f"Extracting JSON: {path}")
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return json.dumps(json.load(f), indent=2)


def extract_yaml(path):
    logger.info(f"Extracting YAML: {path}")
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return yaml.dump(yaml.safe_load(f), sort_keys=False)


def extract_csv(path):
    logger.info(f"Extracting CSV: {path}")
    df = pd.read_csv(path)
    return df.to_string(index=False)


def extract_excel(path):
    logger.info(f"Extracting Excel: {path}")
    excel = pd.ExcelFile(path)

    output = []
    for sheet in excel.sheet_names:
        df = pd.read_excel(path, sheet_name=sheet)
        output.append(f"\n=== SHEET: {sheet} ===\n")
        output.append(df.to_string(index=False))

    return "\n".join(output)


def extract_docx(path):
    logger.info(f"Extracting DOCX: {path}")
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_pptx(path):
    logger.info(f"Extracting PPTX: {path}")
    prs = Presentation(path)

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slides.append(f"\n=== SLIDE {i} ===\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides.append(shape.text)

    return "\n".join(slides)


def extract_pdf(path):
    logger.info(f"Extracting PDF: {path}")
    reader = PdfReader(path)

    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        pages.append(f"\n=== PAGE {i} ===\n")
        if text:
            pages.append(text)

    return "\n".join(pages)


def extract_image_ocr(path):
    logger.info(f"OCR Image: {path}")
    image = Image.open(path)
    return pytesseract.image_to_string(image)


# =========================================================
# ROUTER
# =========================================================

def extract_content(path):

    ext = Path(path).suffix.lower()

    if ext not in CONFIG["SUPPORTED_EXTENSIONS"]:
        raise Exception(f"Unsupported file type: {ext}")

    if ext in [".txt", ".md"]:
        return extract_text_file(path)

    if ext == ".json":
        return extract_json(path)

    if ext in [".yaml", ".yml"]:
        return extract_yaml(path)

    if ext == ".csv":
        return extract_csv(path)

    if ext == ".xlsx":
        return extract_excel(path)

    if ext == ".docx":
        return extract_docx(path)

    if ext == ".pptx":
        return extract_pptx(path)

    if ext == ".pdf":
        return extract_pdf(path)

    if ext in [".png", ".jpg", ".jpeg"]:
        return extract_image_ocr(path)

    raise Exception("No extractor found")


# =========================================================
# NORMALIZER
# =========================================================

def normalize_content(item):

    logger.info(f"Processing input: {item}")

    result = {
        "source_type": "",
        "source_name": "",
        "content": "",
        "metadata": {},
        "truncated": False,
        "status": "SUCCESS"
    }

    input_type = classify_input(item)
    logger.info(f"Classified as: {input_type}")

    # -----------------------------------------------------
    # RAW TEXT
    # -----------------------------------------------------
    if input_type == "RAW_TEXT":

        text = normalize_text(str(item))
        text, truncated = truncate_text(text, CONFIG["MAX_CONTENT_CHARS_PER_ITEM"])

        result.update({
            "source_type": "RAW_TEXT",
            "source_name": "USER_INPUT",
            "content": text,
            "truncated": truncated
        })

        return result

    # -----------------------------------------------------
    # FILE
    # -----------------------------------------------------
    path = Path(item)

    result["source_type"] = path.suffix.lower().replace(".", "").upper()
    result["source_name"] = path.name

    try:

        if not size_limit_check(item):
            logger.warning(f"File too large, skipping: {item}")
            result["status"] = "SKIPPED"
            result["content"] = "[FILE SIZE LIMIT EXCEEDED]"
            return result

        extracted = extract_content(item)

        logger.info(f"Extracted {len(extracted)} chars from {item}")

        text = normalize_text(extracted)

        text, truncated = truncate_text(
            text,
            CONFIG["MAX_CONTENT_CHARS_PER_ITEM"]
        )

        result["content"] = text
        result["truncated"] = truncated

        if truncated:
            logger.warning(f"Content truncated: {item}")

        return result

    except Exception as e:
        logger.error(f"Extraction failed for {item}: {e}")

        result["status"] = "ERROR"
        result["content"] = f"[ERROR: {e}]"

        return result


# =========================================================
# FINAL CONTEXT BUILDER
# =========================================================

def build_final_context(items):

    logger.info("Building final context")

    sections = []

    for item in items:

        sections.append(f"""
========================================
SOURCE: {item['source_name']}
TYPE: {item['source_type']}
STATUS: {item['status']}
========================================

{item['content']}
""")

    final = "\n".join(sections)

    if len(final) > CONFIG["MAX_FINAL_CONTEXT_CHARS"]:
        logger.warning("Final context truncated")
        final = final[:CONFIG["MAX_FINAL_CONTEXT_CHARS"]]

    return final


# =========================================================
# MAIN ENTRY
# =========================================================

def process_inputs(inputs):

    logger.info("Starting universal normalization pipeline")

    if isinstance(inputs, str):
        inputs = [inputs]

    if len(inputs) > CONFIG["MAX_INPUTS"]:
        logger.warning("Input count exceeded limit, truncating")
        inputs = inputs[:CONFIG["MAX_INPUTS"]]

    normalized = []

    for item in inputs:
        normalized.append(normalize_content(item))

    final_context = build_final_context(normalized)

    logger.info("Normalization complete")

    return {
        "status": "SUCCESS",
        "summary": {
            "total_inputs": len(inputs),
            "processed": len(normalized)
        },
        "normalized_inputs": normalized,
        "final_context": final_context
    }


# =========================================================
# EXAMPLE USAGE
# =========================================================

if __name__ == "__main__":

    sample_inputs = [
        "Build executive reporting system",
        "sample.txt",
        "roadmap.csv",
        "presentation.pptx",
        "design.docx",
        "metrics.xlsx",
        "notes.md",
        "diagram.png",
        "architecture.pdf"
    ]

    result = process_inputs(sample_inputs)

    print(json.dumps(result["summary"], indent=2))
    print("\n===== FINAL CONTEXT =====\n")
    print(result["final_context"])